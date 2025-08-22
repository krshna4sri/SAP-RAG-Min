"""
SAP RAG — minimal, production-minded starter
==========================================

What this gives you (in one file):
- Ingest SAP docs (PDF/DOCX/PPTX/TXT/MD) and BPMN XML exports (Signavio/ARIS/SolMan) with page/slide-aware metadata.
- Chunking tuned for procedures/policies with overlap.
- Hybrid retrieval: dense (embeddings) + optional BM25 keyword match.
- Optional cross-encoder re-ranking for quality.
- Q&A generation via:
  1) OpenAI (if OPENAI_API_KEY is set), or
  2) Local Ollama (if OLLAMA_MODEL is set and server is running), or
  3) Fallback extractive QA (Roberta squad2) when no LLM is available.
- FastAPI endpoints (/ingest, /query) and a tiny CLI.
- Source citing (doc path + page/slide + text preview) and confidence score.

Quickstart
----------
1) Python 3.10+
2) Install deps (GPU optional, works on CPU):
   pip install -U fastapi uvicorn[standard] pydantic
   pip install -U sentence-transformers faiss-cpu numpy rank-bm25
   pip install -U pypdf python-docx python-pptx lxml
   pip install -U transformers torch --extra-index-url https://download.pytorch.org/whl/cpu
   # Optional extras for better re-ranking and OpenAI/Ollama support
   pip install -U openai httpx

3) Start the API:
   uvicorn sap_rag_minimal:app --reload --port 8000

4) Ingest your folder of documents (once):
   curl -X POST "http://localhost:8000/ingest?root=/absolute/path/to/your/docs"

5) Ask a question:
   curl "http://localhost:8000/query?q=How do we post an incoming invoice in SAP S/4HANA?"

Environment knobs
-----------------
- EMBEDDING_MODEL (default: "BAAI/bge-small-en-v1.5")
- CROSS_ENCODER (default: "cross-encoder/ms-marco-MiniLM-L-6-v2")
- INDEX_DIR (default: ".rag_index")
- OPENAI_API_KEY (optional): enables OpenAI for answer generation
- OPENAI_MODEL (default: "gpt-4o-mini")
- OLLAMA_MODEL (optional, e.g. "llama3.1") and OLLAMA_BASE_URL (default: http://localhost:11434)

Notes for SAP/BPM teams
-----------------------
- Export process maps as BPMN XML to capture task/gateway names; this file is supported.
- T-codes and module synonyms are auto-expanded (see SAP_SYNONYMS).
- For PDFs, we keep page numbers and preview snippets to aid validation.
- For production, consider pgvector, Qdrant, or Milvus instead of FAISS-on-disk.
- Replace OpenAI with your preferred provider easily in _generate_answer().

"""
from __future__ import annotations
import os, re, json, math, uuid, time, pathlib, logging
from dataclasses import dataclass, asdict
from typing import List, Dict, Any, Optional, Tuple

import numpy as np

# Vector & NLP
import faiss
from sentence_transformers import SentenceTransformer

# Keyword retrieval
from rank_bm25 import BM25Okapi

# File parsers
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from lxml import etree

# Optional LLMs
import httpx

# API
from fastapi import FastAPI, Query
from pydantic import BaseModel

logging.basicConfig(level=logging.INFO)
log = logging.getLogger("sap_rag")

# -----------------------------
# Domain helpers (SAP synonyms)
# -----------------------------
SAP_SYNONYMS: Dict[str, List[str]] = {
    "S/4HANA": ["S4HANA", "SAP S4", "SAP S/4"],
    "FI": ["Finance", "Financial Accounting"],
    "CO": ["Controlling"],
    "MM": ["Materials Management"],
    "SD": ["Sales and Distribution"],
    "PP": ["Production Planning"],
    "HCM": ["HR", "Human Capital Management"],
    "FB60": ["Vendor Invoice", "Post Vendor Invoice"],
    "MIRO": ["Logistics Invoice Verification", "Invoice Verification"],
    "FB50": ["G/L document posting"],
    "MIGO": ["Goods Movement", "Goods Receipt"],
}

# -----------------------------
# Data structures
# -----------------------------
@dataclass
class Chunk:
    id: str
    doc_id: str
    source_path: str
    location: str  # e.g., "page 3" or "slide 7" or "bpmn"
    text: str

@dataclass
class IndexStore:
    # Dense index
    faiss_index: Any
    embeddings: np.ndarray  # (N, D)
    chunks: List[Chunk]
    # Keyword index
    bm25: Optional[BM25Okapi]
    tokenized_corpus: Optional[List[List[str]]]

# -----------------------------
# Utilities
# -----------------------------

def normalize_ws(s: str) -> str:
    s = re.sub(r"\s+", " ", s)
    return s.strip()

def split_into_chunks(text: str, chunk_words: int = 250, overlap_words: int = 50) -> List[str]:
    words = text.split()
    if not words:
        return []
    chunks = []
    step = max(1, chunk_words - overlap_words)
    for start in range(0, len(words), step):
        part = words[start:start+chunk_words]
        if not part:
            break
        chunks.append(" ".join(part))
        if start + chunk_words >= len(words):
            break
    return chunks

# -----------------------------
# Parsers
# -----------------------------

def parse_pdf(path: str) -> List[Tuple[str, str]]:
    out = []
    reader = PdfReader(path)
    for i, page in enumerate(reader.pages, start=1):
        try:
            txt = page.extract_text() or ""
        except Exception:
            txt = ""
        txt = normalize_ws(txt)
        if txt:
            out.append((f"page {i}", txt))
    return out


def parse_docx(path: str) -> List[Tuple[str, str]]:
    doc = DocxDocument(path)
    text = []
    for p in doc.paragraphs:
        if p.text:
            text.append(p.text)
    full = normalize_ws("\n".join(text))
    return [("docx", t) for t in split_into_chunks(full, 300, 60)]


def parse_pptx(path: str) -> List[Tuple[str, str]]:
    pres = Presentation(path)
    out = []
    for i, slide in enumerate(pres.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    texts.append(p.text)
        txt = normalize_ws(" ".join(texts))
        if txt:
            out.append((f"slide {i}", txt))
    return out


def parse_txt_md(path: str) -> List[Tuple[str, str]]:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        txt = normalize_ws(f.read())
    return [("text", t) for t in split_into_chunks(txt, 300, 60)]


def parse_bpmn_xml(path: str) -> List[Tuple[str, str]]:
    """Very light BPMN reader: collects element names to create a linearized process narrative."""
    try:
        tree = etree.parse(path)
    except Exception:
        return []
    root = tree.getroot()
    ns = {"bpmn": "http://www.omg.org/spec/BPMN/20100524/MODEL"}
    labels = []
    for tag in [
        "process", "startEvent", "task", "userTask", "serviceTask", "exclusiveGateway",
        "parallelGateway", "inclusiveGateway", "endEvent", "subProcess", "callActivity"
    ]:
        for node in root.findall(f".//bpmn:{tag}", ns):
            name = node.get("name")
            if name:
                labels.append(name)
    text = " -> ".join(labels)
    text = normalize_ws(text)
    if not text:
        return []
    return [("bpmn", t) for t in split_into_chunks(text, 200, 40)]


PARSERS = {
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".pptx": parse_pptx,
    ".txt": parse_txt_md,
    ".md": parse_txt_md,
    ".bpmn": parse_bpmn_xml,
    ".xml": parse_bpmn_xml,  # treat unknown xml as BPMN attempt
}

# -----------------------------
# Indexer
# -----------------------------
class RagEngine:
    def __init__(self,
                 index_dir: str = os.getenv("INDEX_DIR", ".rag_index"),
                 embed_model_name: str = os.getenv("EMBEDDING_MODEL", "BAAI/bge-small-en-v1.5"),
                 cross_encoder_name: str = os.getenv("CROSS_ENCODER", "cross-encoder/ms-marco-MiniLM-L-6-v2")):
        self.index_dir = index_dir
        os.makedirs(self.index_dir, exist_ok=True)
        log.info(f"Loading embedding model: {embed_model_name}")
        self.embedder = SentenceTransformer(embed_model_name)
        self.embed_dim = self.embedder.get_sentence_embedding_dimension()
        self.cross_encoder_name = cross_encoder_name
        self._index: Optional[IndexStore] = None

    # ---------- Persistence ----------
    def _meta_path(self) -> str:
        return os.path.join(self.index_dir, "meta.json")

    def _faiss_path(self) -> str:
        return os.path.join(self.index_dir, "index.faiss")

    def save(self):
        if not self._index:
            return
        faiss.write_index(self._index.faiss_index, self._faiss_path())
        meta = {
            "chunks": [asdict(c) for c in self._index.chunks],
            # corpus for bm25
            "corpus": self._index.tokenized_corpus,
        }
        with open(self._meta_path(), "w", encoding="utf-8") as f:
            json.dump(meta, f)
        log.info(f"Saved index to {self.index_dir}")

    def load(self):
        if not (os.path.exists(self._faiss_path()) and os.path.exists(self._meta_path())):
            log.warning("No persisted index found; build with ingest().")
            return
        index = faiss.read_index(self._faiss_path())
        with open(self._meta_path(), "r", encoding="utf-8") as f:
            meta = json.load(f)
        chunks = [Chunk(**c) for c in meta["chunks"]]
        tokenized_corpus = meta.get("corpus")
        bm25 = None
        if tokenized_corpus:
            bm25 = BM25Okapi(tokenized_corpus)
        # embeddings are stored inside FAISS only (to save RAM). We'll query via FAISS directly.
        self._index = IndexStore(index, None, chunks, bm25, tokenized_corpus)
        log.info("Index loaded.")

    # ---------- Ingestion ----------
    def ingest(self, root: str, rebuild: bool = True):
        """Walk a folder, parse files, chunk, embed, build hybrid index."""
        all_chunks: List[Chunk] = []
        for dirpath, _, filenames in os.walk(root):
            for fn in filenames:
                ext = pathlib.Path(fn).suffix.lower()
                if ext not in PARSERS:
                    continue
                full = os.path.join(dirpath, fn)
                try:
                    items = PARSERS[ext](full)
                except Exception as e:
                    log.warning(f"Parse failed for {full}: {e}")
                    continue
                doc_id = str(uuid.uuid4())
                for loc, text in items:
                    for piece in split_into_chunks(text, 250, 50) if loc.startswith("page ") else [text]:
                        all_chunks.append(Chunk(
                            id=str(uuid.uuid4()),
                            doc_id=doc_id,
                            source_path=full,
                            location=loc,
                            text=piece,
                        ))
        if not all_chunks:
            raise RuntimeError("No supported documents found to ingest.")

        # Dense embeddings
        texts = [c.text for c in all_chunks]
        log.info(f"Embedding {len(texts)} chunks …")
        embs = self.embedder.encode(texts, batch_size=64, convert_to_numpy=True, show_progress_bar=True, normalize_embeddings=True)
        index = faiss.IndexFlatIP(self.embed_dim)
        index.add(embs)

        # Keyword index
        tokenized = [t.lower().split() for t in texts]
        bm25 = BM25Okapi(tokenized)

        self._index = IndexStore(index, embs, all_chunks, bm25, tokenized)
        self.save()

    # ---------- Querying ----------
    def _expand_query(self, q: str) -> str:
        q_expanded = [q]
        for key, syns in SAP_SYNONYMS.items():
            if key.lower() in q.lower() or any(s.lower() in q.lower() for s in syns):
                q_expanded.extend([key] + syns)
        return " ".join(dict.fromkeys(q_expanded))  # de-dupe order

    def _retrieve(self, q: str, k: int = 8, use_bm25: bool = True) -> List[Tuple[int, float]]:
        if not self._index:
            raise RuntimeError("Index not built. Ingest documents first.")
        # Dense
        q_emb = self.embedder.encode([q], convert_to_numpy=True, normalize_embeddings=True)
        scores, idxs = self._index.faiss_index.search(q_emb, k)
        dense = [(int(i), float(s)) for i, s in zip(idxs[0], scores[0]) if i != -1]

        # BM25
        if use_bm25 and self._index.bm25 is not None:
            tokenized_q = q.lower().split()
            bm25_scores = self._index.bm25.get_scores(tokenized_q)
            top_bm25 = np.argsort(-bm25_scores)[:k]
            keyword = [(int(i), float(bm25_scores[i]/(np.max(bm25_scores)+1e-9))) for i in top_bm25]
        else:
            keyword = []

        # Merge (simple max score per id)
        merged: Dict[int, float] = {}
        for i, s in dense + keyword:
            merged[i] = max(merged.get(i, 0.0), s)
        # Sort by score
        ranked = sorted(merged.items(), key=lambda t: t[1], reverse=True)[:k]
        return ranked

    def _rerank(self, q: str, candidates: List[Tuple[int, float]], use_reranker: bool = True) -> List[Tuple[int, float]]:
        if not use_reranker:
            return candidates
        try:
            from transformers import AutoModelForSequenceClassification, AutoTokenizer
            import torch
            tok = AutoTokenizer.from_pretrained(self.cross_encoder_name)
            model = AutoModelForSequenceClassification.from_pretrained(self.cross_encoder_name)
            inputs = tok([q]*len(candidates), [self._index.chunks[i].text for i,_ in candidates], padding=True, truncation=True, return_tensors="pt")
            with torch.no_grad():
                scores = model(**inputs).logits.squeeze(-1).tolist()
            rescored = [(candidates[i][0], float(scores[i])) for i in range(len(candidates))]
            return sorted(rescored, key=lambda t: t[1], reverse=True)
        except Exception as e:
            log.warning(f"Reranker unavailable, skipping: {e}")
            return candidates

    def _generate_answer(self, q: str, contexts: List[Chunk]) -> Tuple[str, float]:
        """Return (answer, confidence). Tries OpenAI, then Ollama, then extractive QA fallback."""
        context_text = "\n\n".join([c.text for c in contexts])[:12000]
        system = (
            "You are a careful SAP/Business Process assistant."
            " Answer ONLY from the provided context. If unsure, say you don't know."
            " Use numbered citations like [1], [2] that map to sources provided by the tool."
        )
        prompt = f"Question: {q}\n\nContext:\n{context_text}\n\nAnswer succinctly with citations."

        # 1) OpenAI
        api_key = os.getenv("OPENAI_API_KEY")
        model = os.getenv("OPENAI_MODEL", "gpt-4o-mini")
        if api_key:
            try:
                from openai import OpenAI
                client = OpenAI(api_key=api_key)
                resp = client.chat.completions.create(
                    model=model,
                    messages=[{"role": "system", "content": system}, {"role": "user", "content": prompt}],
                    temperature=0.2,
                )
                ans = resp.choices[0].message.content.strip()
                # Heuristic confidence: more tokens retrieved -> higher
                conf = min(0.95, 0.6 + 0.05*len(contexts))
                return ans, conf
            except Exception as e:
                log.warning(f"OpenAI generation failed: {e}")

        # 2) Ollama
        ollama_model = os.getenv("OLLAMA_MODEL")
        if ollama_model:
            try:
                base = os.getenv("OLLAMA_BASE_URL", "http://localhost:11434")
                data = {"model": ollama_model, "messages": [
                    {"role": "system", "content": system},
                    {"role": "user", "content": prompt}
                ], "stream": False}
                r = httpx.post(f"{base}/v1/chat/completions", json=data, timeout=120)
                r.raise_for_status()
                j = r.json()
                ans = j["choices"][0]["message"]["content"].strip()
                conf = min(0.9, 0.55 + 0.05*len(contexts))
                return ans, conf
            except Exception as e:
                log.warning(f"Ollama generation failed: {e}")

        # 3) Extractive QA fallback
        try:
            from transformers import AutoTokenizer, AutoModelForQuestionAnswering
            import torch
            model_name = "deepset/roberta-base-squad2"
            tok = AutoTokenizer.from_pretrained(model_name)
            m = AutoModelForQuestionAnswering.from_pretrained(model_name)
            inputs = tok(q, context_text, return_tensors="pt", truncation=True, max_length=512)
            with torch.no_grad():
                out = m(**inputs)
            start = int(out.start_logits.argmax())
            end = int(out.end_logits.argmax())
            ans = tok.decode(inputs["input_ids"][0][start:end+1])
            if not ans.strip():
                ans = "I couldn't find a precise answer in the provided documents."
            conf = 0.5
            return ans, conf
        except Exception as e:
            log.warning(f"Extractive QA failed: {e}")
            return "I'm unable to generate an answer right now.", 0.0

    def ask(self, q: str, top_k: int = 6, use_reranker: bool = True) -> Dict[str, Any]:
        q_expanded = self._expand_query(q)
        cand = self._retrieve(q_expanded, k=top_k*2)
        cand = self._rerank(q, cand, use_reranker)
        # Take top_k
        cand = cand[:top_k]
        chunks = [self._index.chunks[i] for i,_ in cand]
        answer, conf = self._generate_answer(q, chunks)
        sources = []
        for rank, (i, score) in enumerate(cand, start=1):
            ch = self._index.chunks[i]
            preview = (ch.text[:240] + "…") if len(ch.text) > 240 else ch.text
            sources.append({
                "rank": rank,
                "doc_path": ch.source_path,
                "location": ch.location,
                "score": round(float(score), 4),
                "preview": preview,
                "citation": f"[{rank}]",
            })
        return {
            "question": q,
            "answer": answer,
            "confidence": round(conf, 3),
            "sources": sources,
        }

# -----------------------------
# FastAPI
# -----------------------------
app = FastAPI(title="SAP RAG Minimal", version="0.1.0")
ENGINE = RagEngine()
ENGINE.load()

class QueryResp(BaseModel):
    question: str
    answer: str
    confidence: float
    sources: List[Dict[str, Any]]

@app.post("/ingest")
def ingest(root: str = Query(..., description="Root folder to crawl for documents")):
    ENGINE.ingest(root)
    return {"status": "ok", "index_dir": ENGINE.index_dir}

@app.get("/query", response_model=QueryResp)
def query(q: str, top_k: int = 6):
    return ENGINE.ask(q, top_k=top_k)

# -----------------------------
# CLI for quick testing
# -----------------------------
if __name__ == "__main__":
    import argparse
    p = argparse.ArgumentParser(description="SAP RAG minimal CLI")
    sub = p.add_subparsers(dest="cmd")

    p_ing = sub.add_parser("ingest")
    p_ing.add_argument("root", help="Root folder of documents")

    p_q = sub.add_parser("ask")
    p_q.add_argument("question", help="Question to ask")
    p_q.add_argument("--k", type=int, default=6)

    args = p.parse_args()

    eng = RagEngine()
    eng.load()

    if args.cmd == "ingest":
        eng.ingest(args.root)
    elif args.cmd == "ask":
        if not eng._index:
            log.error("Index not found. Run ingest first.")
            raise SystemExit(1)
        out = eng.ask(args.question, top_k=args.k)
        print(json.dumps(out, indent=2))
    else:
        p.print_help()
